"""Build the Monarch defence deck as an editable .pptx.

Diagrams are native shapes rather than images so the presenter can move or retype
anything on the day. Figures that come from the analysis stay as images, because
they are evidence and must not be redrawn by hand.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ICONS = os.path.join(HERE, "icon_png")
REPO = r"C:\Users\Windows\Downloads\monarch"
FIG = os.path.join(REPO, "services", "inference", "data", "figures")
REPORT = os.path.join(REPO, "services", "inference", "data", "final", "report")
PAPER1 = os.path.join(REPO, "services", "inference", "data", "paper1")
OUT = os.path.join(REPO, "docs", "defence", "Monarch_Defence.pptx")

NAVY = RGBColor(0x1F, 0x4E, 0x79)
MIDBLUE = RGBColor(0x2E, 0x75, 0xB6)
EMBER = RGBColor(0xC0, 0x39, 0x2B)
SLATE = RGBColor(0x59, 0x59, 0x59)
MIST = RGBColor(0xD9, 0xD9, 0xD9)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
W, H = 13.333, 7.5
M = 0.75          # slide margin
TITLE_TOP = 0.42
BODY_TOP = 1.42

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size=20, bold=False, colour=BLACK, space_after=8,
         first=False, align=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = FONT
    return p


def rich(tf, parts, size=20, space_after=8, first=False, align=None):
    """One paragraph, several runs: parts is a list of (text, bold, colour)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    for text, bold, colour in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = FONT
    return p


def formula(tf, parts, size=30, colour=NAVY, first=False, align=PP_ALIGN.CENTER):
    """A display formula. Parts are (text, level) with level 0 normal,
    -1 subscript, +1 superscript, so a subscript is set rather than faked with
    an underscore."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(0)
    for text, level in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = colour
        run.font.name = FONT
        if level < 0:
            run.font._rPr.set("baseline", "-25000")
        elif level > 0:
            run.font._rPr.set("baseline", "30000")
    return p


def action_title(s, text, icon=None):
    x = M
    if icon:
        s.shapes.add_picture(os.path.join(ICONS, icon + ".png"),
                             Inches(M), Inches(TITLE_TOP - 0.02),
                             height=Inches(0.46))
        x = M + 0.66
    tf = textbox(s, x, TITLE_TOP - 0.08, W - x - M, 0.9)
    para(tf, text, size=26, bold=True, colour=NAVY, first=True, space_after=0)


def footer(s, section, number):
    tf = textbox(s, M, H - 0.62, 5.0, 0.35)
    para(tf, section, size=11, colour=SLATE, first=True, space_after=0)
    tf = textbox(s, W - M - 1.0, H - 0.62, 1.0, 0.35, align=PP_ALIGN.RIGHT)
    para(tf, str(number), size=11, colour=SLATE, first=True, space_after=0)


def source(s, text, y=None):
    tf = textbox(s, M, y if y else H - 1.05, W - 2 * M, 0.42)
    para(tf, text, size=11, colour=SLATE, first=True, space_after=0)


def card(s, x, y, w, h, line=MIST, fill=WHITE, width=1.25):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(width)
    shp.shadow.inherit = False
    shp.text_frame.word_wrap = True
    return shp


def rule(s, x, y, w, colour=MIST, thickness=1.5):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Pt(thickness))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def arrow(s, x, y, w, h=0.16, colour=SLATE):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def divider(s, num, title, sub, icon):
    s.shapes.add_picture(os.path.join(ICONS, icon + ".png"),
                         Inches(1.3), Inches(2.35), height=Inches(0.72))
    tf = textbox(s, 1.3, 3.25, 9.0, 2.0)
    para(tf, num, size=40, bold=True, colour=MIDBLUE, first=True, space_after=2)
    para(tf, title, size=30, bold=True, colour=NAVY, space_after=8)
    para(tf, sub, size=17, colour=SLATE)


def hero(s, label, number, lines):
    tf = textbox(s, 1.0, 2.35, W - 2.0, 0.6, align=PP_ALIGN.CENTER)
    para(tf, label.upper(), size=15, colour=SLATE, first=True,
         align=PP_ALIGN.CENTER, space_after=0)
    tf = textbox(s, 1.0, 2.95, W - 2.0, 1.5, align=PP_ALIGN.CENTER)
    para(tf, number, size=72, bold=True, colour=NAVY, first=True,
         align=PP_ALIGN.CENTER, space_after=0)
    tf = textbox(s, 2.0, 4.55, W - 4.0, 1.4, align=PP_ALIGN.CENTER)
    for i, (text, colour) in enumerate(lines):
        para(tf, text, size=18, colour=colour, first=(i == 0),
             align=PP_ALIGN.CENTER, space_after=6)


n = 0


def num():
    global n
    n += 1
    return n


# ----------------------------------------------------------------- 1 title
s = slide()
tf = textbox(s, M + 0.15, 1.35, W - 2 * M, 0.4)
para(tf, "B.Sc. PHYSICS   ·   DEPARTMENT OF NATURAL SCIENCES   ·   CUEA",
     size=14, colour=SLATE, first=True, space_after=0)
tf = textbox(s, M + 0.15, 1.95, 11.0, 1.6)
para(tf, "Measuring the External Field", size=44, bold=True, colour=NAVY,
     first=True, space_after=6)
tf = textbox(s, M + 0.15, 3.15, 10.0, 1.0)
para(tf, "A cortical-proxy content observable and the mean-field bound on "
         "media-driven opinion change", size=19, colour=SLATE, first=True,
     space_after=0)
rule(s, M + 0.15, 4.45, 11.0)
tf = textbox(s, M + 0.15, 4.70, 6.0, 0.9)
rich(tf, [("Brian Mwai", True, BLACK), ("   1050555", False, SLATE)],
     size=19, first=True, space_after=4)
para(tf, "3 September 2026", size=14, colour=SLATE, space_after=0)
tf = textbox(s, 7.6, 4.70, 5.0, 0.6, align=PP_ALIGN.RIGHT)
para(tf, "Supervisor: Dr. Songa Mutambi", size=15, colour=SLATE, first=True,
     space_after=0)
num()

# ----------------------------------------------------------------- 2 the claim
s = slide()
s.shapes.add_picture(os.path.join(ICONS, "waveform.png"), Inches(M + 0.15),
                     Inches(1.55), height=Inches(0.5))
tf = textbox(s, M + 0.15, 2.25, 10.6, 1.4)
para(tf, "Statistical physics has always modelled media as a field pushing on a "
         "population of coupled opinions.", size=22, colour=BLACK, first=True,
     space_after=10)
para(tf, "That field is assumed. It is not measured.", size=22, colour=BLACK,
     space_after=0)
rule(s, M + 0.15, 4.05, 10.6)
tf = textbox(s, M + 0.15, 4.35, 10.6, 1.6)
para(tf, "This project measures one from content itself, applies it to 400 "
         "articles, and derives how strong the coupling would have to be for a "
         "spread that size to move a population at all.",
     size=22, bold=True, colour=NAVY, first=True, space_after=0)
num()

# ----------------------------------------------------------------- 3 divider
s = slide()
divider(s, "01", "The problem", "What is missing from the way media is modelled.",
        "newspaper")
num()

# ----------------------------------------------------------------- 4 two articles
s = slide()
action_title(s, "Two articles can report one event and leave readers in "
                "different states", "newspaper")
c1 = card(s, M, 1.75, 5.5, 1.85)
tf = c1.text_frame
tf.margin_left, tf.margin_top = Inches(0.22), Inches(0.16)
para(tf, "WIRE REPORT", size=12, colour=SLATE, first=True, space_after=6)
para(tf, "Federal Reserve holds interest rates steady, citing a stable "
         "inflation outlook.", size=18, colour=BLACK, space_after=0)
c2 = card(s, 7.1, 1.75, 5.5, 1.85)
tf = c2.text_frame
tf.margin_left, tf.margin_top = Inches(0.22), Inches(0.16)
para(tf, "THE SAME EVENT", size=12, colour=SLATE, first=True, space_after=6)
rich(tf, [("FED DESTROYS AMERICA. ", True, BLACK),
          ("Your savings are GONE. The collapse they hid from you!", False, BLACK)],
     size=18, space_after=0)
tf = textbox(s, M, 3.72, 5.5, 0.4, align=PP_ALIGN.CENTER)
para(tf, "informs", size=17, colour=SLATE, first=True, align=PP_ALIGN.CENTER,
     space_after=0)
tf = textbox(s, 7.1, 3.72, 5.5, 0.4, align=PP_ALIGN.CENTER)
para(tf, "provokes", size=17, bold=True, colour=EMBER, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 4.55, W - 2 * M, 1.2)
para(tf, "Readers notice the difference and cannot quantify it. The platforms "
         "distributing both have no objective measure of which they are "
         "amplifying.", size=20, colour=BLACK, first=True, space_after=0)
footer(s, "The problem", num())

# ----------------------------------------------------------------- 5 words vs response
s = slide()
action_title(s, "Existing tools measure the words; none measures the predicted "
                "response", "text-aa")
c1 = card(s, 1.35, 2.15, 4.3, 1.5)
tf = c1.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "SENTIMENT, READABILITY, STANCE", size=12, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=6)
rich(tf, [("characterise ", False, BLACK), ("the words", True, BLACK)],
     size=19, align=PP_ALIGN.CENTER, space_after=0)
arrow(s, 6.05, 2.80, 1.2, 0.2, MIDBLUE)
c2 = card(s, 7.65, 2.15, 4.3, 1.5, line=NAVY, width=2.0)
tf = c2.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "THIS PROJECT", size=12, colour=NAVY, first=True,
     align=PP_ALIGN.CENTER, space_after=6)
rich(tf, [("estimate ", False, BLACK), ("the response", True, BLACK)],
     size=19, align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 4.35, W - 2 * M, 1.6)
para(tf, "A released neural encoder predicts cortical activity from the "
         "content, and the measurement is taken on that prediction.",
     size=20, colour=BLACK, first=True, space_after=10)
para(tf, "The instrument is a thermometer for one property of media. The thesis "
         "is an account of building it, running it, and stating what its "
         "readings support.", size=18, colour=SLATE, space_after=0)
footer(s, "The problem", num())

# ----------------------------------------------------------------- 6 divider
s = slide()
divider(s, "02", "The instrument", "From a passage of text to a single number.",
        "brain")
num()

# ----------------------------------------------------------------- 7 cascade
s = slide()
action_title(s, "A five-stage cascade turns one article into one number", "brain")
stages = [("Text", "one article"), ("Speech", "synthesised"),
          ("Timings", "per word"), ("Embeddings", "text + audio"),
          ("TRIBE v2", "encoder"), ("20,484", "vertices")]
x0, cw, gap, cy, ch = 0.62, 1.72, 0.34, 1.70, 1.05
for i, (head, sub) in enumerate(stages):
    x = x0 + i * (cw + gap)
    hot = i >= 4
    shp = card(s, x, cy, cw, ch, line=NAVY if hot else MIST,
               width=2.0 if hot else 1.25)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, head, size=15, bold=True, colour=BLACK, first=True,
         align=PP_ALIGN.CENTER, space_after=2)
    para(tf, sub, size=12, colour=SLATE, align=PP_ALIGN.CENTER, space_after=0)
    if i < len(stages) - 1:
        arrow(s, x + cw + 0.05, cy + ch / 2 - 0.07, gap - 0.10, 0.14)

x_last = x0 + 5 * (cw + gap)
vline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(x_last + cw / 2 - 0.01), Inches(cy + ch),
                           Pt(1.5), Inches(0.45))
vline.fill.solid()
vline.fill.fore_color.rgb = SLATE
vline.line.fill.background()
vline.shadow.inherit = False
hline = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.05),
                           Inches(cy + ch + 0.45), Inches(x_last + cw / 2 - 2.05),
                           Pt(1.5))
hline.fill.solid()
hline.fill.fore_color.rgb = SLATE
hline.line.fill.background()
hline.shadow.inherit = False

reduce_y = 3.85
for x, head, sub in [(1.45, "Affective", "1,030 vertices"),
                     (5.20, "Deliberative", "851 vertices")]:
    drop = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x + 1.16),
                              Inches(cy + ch + 0.45), Inches(0.18), Inches(0.42))
    drop.fill.solid()
    drop.fill.fore_color.rgb = SLATE
    drop.line.fill.background()
    drop.shadow.inherit = False
    shp = card(s, x, reduce_y, 2.5, 1.0)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, head, size=15, bold=True, colour=BLACK, first=True,
         align=PP_ALIGN.CENTER, space_after=2)
    para(tf, sub, size=12, colour=SLATE, align=PP_ALIGN.CENTER, space_after=0)

tf = textbox(s, 4.05, reduce_y + 0.28, 1.0, 0.5, align=PP_ALIGN.CENTER)
para(tf, "−", size=26, bold=True, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, 7.80, reduce_y + 0.28, 1.0, 0.5, align=PP_ALIGN.CENTER)
para(tf, "=", size=26, bold=True, colour=NAVY, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
shp = card(s, 8.80, reduce_y, 2.9, 1.0, line=NAVY, width=2.0)
tf = shp.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, "One number", size=15, bold=True, colour=NAVY, first=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "per article", size=12, colour=SLATE, align=PP_ALIGN.CENTER,
     space_after=0)
source(s, "Encoder: TRIBE v2, d'Ascoli et al., Meta FAIR (2025), arXiv:2507.22229.  "
          "Regions from HCP MMP1.0, Glasser et al. (2016).  "
          "Instrument at monarch-4iy.pages.dev", y=5.55)
footer(s, "The instrument", num())

# ----------------------------------------------------------------- 8 regions
s = slide()
action_title(s, "The index contrasts two cortical networks, 1,030 vertices "
                "against 851", "brain")
s.shapes.add_picture(os.path.join(FIG, "B1_roi_definition.png"),
                     Inches(1.55), Inches(1.50), width=Inches(10.2))
tf = textbox(s, M, 5.30, W - 2 * M, 0.9, align=PP_ALIGN.CENTER)
para(tf, "Affective salience in orange, deliberative control in blue. Resolved "
         "once into vertex index sets and cached, so all 400 items are averaged "
         "over identical regions.", size=16, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
source(s, "Regions defined on HCP MMP1.0 (Glasser et al., 2016); surfaces are "
          "fsaverage5.")
footer(s, "The instrument", num())

# ----------------------------------------------------------------- 9 observable
s = slide()
action_title(s, "The observable is a signed difference, defined for every item",
             "function")
tf = textbox(s, M, 2.00, W - 2 * M, 1.1, align=PP_ALIGN.CENTER)
formula(tf, [("NAA", 0), ("signed", -1), ("  =  A", 0), ("affective", -1),
             ("  −  A", 0), ("deliberative", -1)], size=38, first=True)
rule(s, 3.1, 3.75, 7.1)
tick = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.63), Inches(3.60),
                          Pt(1.5), Inches(0.28))
tick.fill.solid()
tick.fill.fore_color.rgb = BLACK
tick.line.fill.background()
tick.shadow.inherit = False
tf = textbox(s, 6.25, 3.90, 0.8, 0.35, align=PP_ALIGN.CENTER)
para(tf, "0", size=14, colour=SLATE, first=True, align=PP_ALIGN.CENTER,
     space_after=0)
tf = textbox(s, 3.1, 3.25, 2.6, 0.4)
para(tf, "reasoning leads", size=17, bold=True, colour=EMBER, first=True,
     space_after=0)
tf = textbox(s, 7.6, 3.25, 2.6, 0.4, align=PP_ALIGN.RIGHT)
para(tf, "emotion leads", size=17, bold=True, colour=NAVY, first=True,
     align=PP_ALIGN.RIGHT, space_after=0)
tf = textbox(s, M, 4.85, W - 2 * M, 1.0)
rich(tf, [("A cortical proxy", True, BLACK),
          (": it rates content, never a person, and is never called the "
           "amygdala.", False, BLACK)], size=20, first=True, space_after=0)
footer(s, "The instrument", num())

# ----------------------------------------------------------------- 10 checkpoint
s = slide()
action_title(s, "Inspecting the checkpoint changed what the project could claim",
             "magnifying-glass")
findings = [
    ("01", "Cortical only",
     "Output is (T, 20484). There is no subcortical head.",
     "The amygdala is not predicted, so the proposal's central equation cannot "
     "be computed.", True),
    ("02", "Averaged over subjects",
     "The loader forces average_subjects.",
     "It predicts a typical viewer. Whether that carries signal is the third "
     "paper's question.", False),
    ("03", "Standardised output",
     "Region means sit near zero, so a ratio sign-flips or explodes.",
     "Undefined for 69 of 400 items. Replaced by the difference.", False),
]
cw = 3.75
for i, (idx, head, body, tail, alert) in enumerate(findings):
    x = M + i * (cw + 0.45)
    tf = textbox(s, x, 1.65, cw, 3.4)
    para(tf, idx, size=20, bold=True, colour=MIDBLUE, first=True, space_after=4)
    para(tf, head, size=19, bold=True, colour=BLACK, space_after=8)
    para(tf, body, size=16, colour=BLACK, space_after=8)
    para(tf, tail, size=16, colour=EMBER if alert else BLACK, space_after=0)
tf = textbox(s, M, 5.35, W - 2 * M, 0.8)
para(tf, "Each was found by testing the artifact, and each is stated in the "
         "abstract rather than the limitations.", size=17, colour=SLATE,
     first=True, space_after=0)
footer(s, "The instrument", num())

# ----------------------------------------------------------------- 11 divider
s = slide()
divider(s, "03", "The measurement", "400 articles, scanned twice.", "chart-bar")
num()

# ----------------------------------------------------------------- 12 corpus
s = slide()
action_title(s, "The corpus is length-matched and was powered before the scan",
             "books")
cats = [("fear-activating", "ISOT-fake"), ("high outrage", "SemEval-2019 T4"),
        ("reward hook", "Webis-Clickbait-17"), ("neutral informational",
                                                "PubMed + ISOT-true")]
cw = 2.85
for i, (name, src) in enumerate(cats):
    x = M + i * (cw + 0.32)
    shp = card(s, x, 1.60, cw, 1.65)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, "100", size=26, bold=True, colour=NAVY, first=True,
         align=PP_ALIGN.CENTER, space_after=4)
    para(tf, name, size=16, bold=True, colour=BLACK, align=PP_ALIGN.CENTER,
         space_after=4)
    para(tf, src, size=12, colour=SLATE, align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 3.55, W - 2 * M, 1.9)
rich(tf, [("Length-matched. ", True, BLACK),
          ("Mean word counts 167.4, 163.5, 163.7, 162.2; standard deviations "
           "near 11.", False, BLACK)], size=20, first=True, space_after=12)
rich(tf, [("Powered first. ", True, BLACK),
          ("Smallest detectable η² = 0.0268, smallest detectable "
           "AUC = 0.5916, both fixed before a single item was scanned.",
           False, BLACK)], size=20, space_after=0)
source(s, "Corpora: Ahmed et al. (ISOT, 2017); Kiesel et al. (SemEval-2019 "
          "Task 4); Potthast et al. (Webis-Clickbait-17); PubMed open abstracts.")
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 13 hero
s = slide()
hero(s, "the four categories separate at", "η² = 0.1068",
     [("F = 15.779    p = 1.03 × 10⁻⁹    n = 400", BLACK),
      ("against a design powered to detect 0.0268; AUC against the pre-scan "
       "label is 0.6274", SLATE)])
num()

# ----------------------------------------------------------------- 14 violin
s = slide()
action_title(s, "Categories shift against each other, but the distributions "
                "overlap", "chart-bar")
s.shapes.add_picture(os.path.join(REPORT, "fig_violin.png"),
                     Inches(3.05), Inches(1.45), height=Inches(4.05))
tf = textbox(s, M, 5.60, W - 2 * M, 0.8, align=PP_ALIGN.CENTER)
para(tf, "Every category mean is negative: the deliberative network leads "
         "throughout, and the categories differ in how far below zero they sit.",
     size=16, colour=SLATE, first=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 15 effect sizes
s = slide()
action_title(s, "Fear-activating content carries the effect; outrage does not "
                "separate", "chart-bar")
axis = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.55), Inches(1.75),
                          Pt(1.5), Inches(2.55))
axis.fill.solid()
axis.fill.fore_color.rgb = MIST
axis.line.fill.background()
axis.shadow.inherit = False
bars = [("fear-activating", 0.939, NAVY, "d = 0.939", None),
        ("reward hook", 0.319, MIDBLUE, "d = 0.319", None),
        ("high outrage", 0.030, EMBER, "d = 0.030", "p = 0.832")]
scale = 7.0
for i, (label, d, colour, dlabel, plabel) in enumerate(bars):
    y = 1.95 + i * 0.85
    tf = textbox(s, M, y - 0.06, 2.65, 0.45, align=PP_ALIGN.RIGHT)
    para(tf, label, size=18, colour=BLACK, first=True, align=PP_ALIGN.RIGHT,
         space_after=0)
    bw = max(d * scale, 0.05)
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.57), Inches(y),
                             Inches(bw), Inches(0.36))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = textbox(s, 3.57 + bw + 0.14, y - 0.06, 3.2, 0.45)
    if plabel:
        rich(tf, [(dlabel + "   ", True, colour), (plabel, False, SLATE)],
             size=17, first=True, space_after=0)
    else:
        para(tf, dlabel, size=17, bold=(i == 0), colour=colour, first=True,
             space_after=0)
tf = textbox(s, M, 4.80, W - 2 * M, 1.4)
para(tf, "Cohen's d against the neutral baseline.", size=20, colour=BLACK,
     first=True, space_after=8)
para(tf, "Outrage is the category the proposal expected to separate most, and "
         "it does not separate at all.", size=20, bold=True, colour=EMBER,
     space_after=0)
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 16 confound
s = slide()
action_title(s, "Category cannot be told apart from source dataset", "warning")
pairs = [("fear-activating", "ISOT-fake"),
         ("high outrage", "SemEval-2019 Task 4"),
         ("reward hook", "Webis-Clickbait-17"),
         ("neutral informational", "PubMed + ISOT-true")]
tf = textbox(s, 1.55, 1.55, 3.4, 0.4, align=PP_ALIGN.CENTER)
para(tf, "CATEGORY", size=13, bold=True, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, 8.35, 1.55, 3.4, 0.4, align=PP_ALIGN.CENTER)
para(tf, "SOURCE DATASET", size=13, bold=True, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
for i, (cat, src) in enumerate(pairs):
    y = 2.05 + i * 0.78
    shp = card(s, 1.55, y, 3.4, 0.6)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, cat, size=15, colour=BLACK, first=True, align=PP_ALIGN.CENTER,
         space_after=0)
    arrow(s, 5.15, y + 0.23, 3.0, 0.14, EMBER)
    shp = card(s, 8.35, y, 3.4, 0.6)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, src, size=15, colour=BLACK, first=True, align=PP_ALIGN.CENTER,
         space_after=0)
tf = textbox(s, M, 5.35, W - 2 * M, 1.2)
para(tf, "One source per category, so the separation is between corpora and "
         "cannot be attributed to framing. Length is controlled; provenance is "
         "not.", size=20, colour=BLACK, first=True, space_after=0)
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 17 reliability
s = slide()
action_title(s, "Group-level claims hold at this precision; per-item claims do "
                "not", "arrows-clockwise")
tf = textbox(s, M, 1.70, 5.6, 1.1, align=PP_ALIGN.CENTER)
para(tf, "0.8725", size=54, bold=True, colour=NAVY, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 2.85, 5.6, 0.45, align=PP_ALIGN.CENTER)
para(tf, "ICC between two GPU sessions", size=16, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 3.45, 5.6, 1.5)
para(tf, "The separation replicates: 0.0888 on the second session against "
         "0.1068 on the first. Both clear the detectable floor.", size=17,
     colour=BLACK, first=True, space_after=0)

tf = textbox(s, 7.1, 1.70, 5.5, 1.1, align=PP_ALIGN.CENTER)
para(tf, "12.8%", size=54, bold=True, colour=EMBER, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, 7.1, 2.85, 5.5, 0.45, align=PP_ALIGN.CENTER)
para(tf, "of items reverse direction", size=16, colour=SLATE, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, 7.1, 3.45, 5.5, 1.5)
para(tf, "51 of 400 items flip sign between sessions, so no verdict on one "
         "article appears anywhere in the analysis, the figures or the public "
         "site.", size=17, colour=BLACK, first=True, space_after=0)
tf = textbox(s, M, 5.35, W - 2 * M, 0.9)
para(tf, "A second session scanned all 400 items with identical text, code and "
         "regions. Neither run is corrected toward the other, and they are not "
         "averaged.", size=17, colour=SLATE, first=True, space_after=0)
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 18 per-vertex
s = slide()
action_title(s, "The encoder returns a full cortical map, not a summary "
                "statistic", "brain")
s.shapes.add_picture(os.path.join(FIG, "B3_per_vertex_fear_activating-0000.png"),
                     Inches(2.35), Inches(1.45), width=Inches(8.6))
tf = textbox(s, M, 5.35, W - 2 * M, 0.9, align=PP_ALIGN.CENTER)
para(tf, "One fear-activating item at all 20,484 vertices, painted above the "
         "map's 70th percentile. Nothing is smoothed, interpolated or simulated.",
     size=16, colour=SLATE, first=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s, "The measurement", num())

# ----------------------------------------------------------------- 19 divider
s = slide()
divider(s, "04", "The physics", "What a measured field is worth.", "atom")
num()

# ----------------------------------------------------------------- 20 mean field
s = slide()
action_title(s, "Media enters the mean-field model as a field on the measured "
                "observable", "atom")
tf = textbox(s, M, 1.60, W - 2 * M, 1.0)
para(tf, "Each person holds an opinion ±1, is pulled toward their neighbours "
         "with strength J, and is pushed by media with field h:", size=19,
     colour=BLACK, first=True, space_after=0)
tf = textbox(s, M, 2.45, W - 2 * M, 0.7, align=PP_ALIGN.CENTER)
para(tf, "m  =  tanh( βJ m + h ),        h  =  α X", size=26, bold=True,
     colour=NAVY, first=True, align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 3.25, W - 2 * M, 0.8)
para(tf, "where X is the measured observable and α the unknown coupling. "
         "Expanding artanh(m) ≈ m + m³/3:", size=19, colour=BLACK, first=True,
     space_after=0)
tf = textbox(s, M, 4.00, W - 2 * M, 0.7, align=PP_ALIGN.CENTER)
para(tf, "F(m) = a m² + b m⁴ − h m,        a = (1 − βJ)/2,        b = 1/12",
     size=24, bold=True, colour=NAVY, first=True, align=PP_ALIGN.CENTER,
     space_after=0)
tf = textbox(s, M, 4.85, W - 2 * M, 1.2)
rich(tf, [("Coefficients derived from the self-consistency condition. ",
           True, SLATE),
          ("The proposal carried a = 1 − βJ, b = (βJ)³/3, which do not follow "
           "from it. The solver reproduces the exact exponents 1/2, 1 and 3.",
           False, SLATE)], size=17, first=True, space_after=0)
footer(s, "The physics", num())

# ----------------------------------------------------------------- 21 no alpha
s = slide()
action_title(s, "The corpus does not identify the coupling, so no value is "
                "quoted", "x-circle")
tf = textbox(s, M, 1.75, 6.6, 3.6)
rich(tf, [("What was tried. ", True, BLACK),
          ("Regress the outcome on the measured index and read off α.",
           False, BLACK)], size=20, first=True, space_after=14)
rich(tf, [("What happened. ", True, BLACK),
          ("The fit performs worse than the sample mean out of sample, and the "
           "estimate scales inversely with a coupling βJ that no measurement "
           "here fixes.", False, BLACK)], size=20, space_after=14)
para(tf, "So no value of α-hat appears anywhere — not in a table, a caption, or "
         "\u201cfor illustration\u201d.", size=20, bold=True, colour=EMBER,
     space_after=0)
tf = textbox(s, 8.05, 1.85, 4.55, 3.0)
para(tf, "Reporting an unidentified coupling as though it were measured would "
         "be the failure here. Declining to quote it is not.", size=17,
     colour=SLATE, first=True, space_after=14)
para(tf, "What replaces it is a statement that does not need the fit at all.",
     size=17, colour=SLATE, space_after=0)
footer(s, "The physics", num())

# ----------------------------------------------------------------- 22 the bound
s = slide()
action_title(s, "Any content observable must clear a bound to drive a "
                "transition", "ruler")
s.shapes.add_picture(os.path.join(PAPER1, "F6_alpha_required.png"),
                     Inches(M), Inches(1.60), height=Inches(3.45))
tf = textbox(s, 7.35, 1.90, 5.25, 0.9, align=PP_ALIGN.CENTER)
formula(tf, [("α  ≥  h", 0), ("c", -1), ("(βJ) / ΔX", 0)], size=32, first=True)
tf = textbox(s, 7.35, 2.95, 5.25, 2.1)
para(tf, "Rather than fitting α and hoping it clears zero, this states what α "
         "would have to be for the mechanism to work at all — for any candidate "
         "observable, before a corpus exists.", size=18, colour=BLACK,
     first=True, space_after=0)
tf = textbox(s, M, 5.35, W - 2 * M, 0.7, align=PP_ALIGN.CENTER)
rich(tf, [("Measured spread ΔX = 0.1241  ⇒  α ≥ 0.169 at βJ = 1.10,   "
           "1.655 at 1.50,   ", False, BLACK), ("4.294 at 2.00", True, NAVY)],
     size=18, first=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s, "The physics", num())

# ----------------------------------------------------------------- 23 divider
s = slide()
divider(s, "05", "Standing", "What holds, and what is still open.", "flag-banner")
num()

# ----------------------------------------------------------------- 24 paper 3
s = slide()
action_title(s, "Whether the encoder tracks real cortex is measurable, and one "
                "pass remains", "target")
tf = textbox(s, M, 1.60, W - 2 * M, 1.0)
para(tf, "A commercial lab reports the released checkpoint, in the averaged "
         "configuration this project runs, is anti-correlated with real cortex. "
         "Four subjects, one stimulus, not independently replicated.", size=19,
     colour=BLACK, first=True, space_after=0)
for x, value, label, note, colour in [
    (1.75, "0.1517", "measured noise ceiling",
     "CI [0.1484, 0.1551], 4 subjects, 1000 parcels", NAVY),
    (7.05, "0.2247", "withdrawn",
     "an episode picked by position, and a guard that could not fire", EMBER)]:
    shp = card(s, x, 2.75, 4.5, 1.55)
    tf = shp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, value, size=30, bold=True, colour=colour, first=True,
         align=PP_ALIGN.CENTER, space_after=2)
    para(tf, label, size=16, colour=BLACK, align=PP_ALIGN.CENTER, space_after=2)
    para(tf, note, size=12, colour=SLATE, align=PP_ALIGN.CENTER, space_after=0)
tf = textbox(s, M, 4.55, W - 2 * M, 1.0)
para(tf, "Replaying the old code path reproduces the wrong number exactly, "
         "which proved the cause. Both defects were silent: nothing crashed.",
     size=19, colour=BLACK, first=True, space_after=0)
source(s, "Schaefer et al. (2018); Algonauts 2025 (CC0); audit self-published, "
          "Sapient (2026).")
footer(s, "Standing", num())

# ----------------------------------------------------------------- 25 claims
s = slide()
action_title(s, "Four claims the work supports, and four it does not",
             "check-circle")
s.shapes.add_picture(os.path.join(ICONS, "check-circle.png"), Inches(M),
                     Inches(1.55), height=Inches(0.34))
tf = textbox(s, M + 0.48, 1.50, 5.0, 0.45)
para(tf, "SUPPORTED", size=17, bold=True, colour=NAVY, first=True, space_after=0)
tf = textbox(s, M, 2.10, 5.6, 3.6)
for text in [
    "A content observable that can be measured, with its reliability stated.",
    "That it separates these four corpora at 0.1068, power attached.",
    "A bound any candidate observable must satisfy, usable before data exists.",
    "A measured noise ceiling for the checkpoint on public data.",
]:
    para(tf, "•  " + text, size=17, colour=BLACK,
         first=(text.startswith("A content")), space_after=12)

s.shapes.add_picture(os.path.join(ICONS, "x-circle.png"), Inches(7.1),
                     Inches(1.55), height=Inches(0.34))
tf = textbox(s, 7.58, 1.50, 5.0, 0.45)
para(tf, "NOT SUPPORTED", size=17, bold=True, colour=EMBER, first=True,
     space_after=0)
tf = textbox(s, 7.1, 2.10, 5.5, 3.6)
for text in [
    "That the index detects manipulation. Category is confounded with source.",
    "Anything about the amygdala. The checkpoint does not predict it.",
    "Any value of the coupling α. The corpus does not identify it.",
    "Any verdict on one article. 12.8% reverse between sessions.",
]:
    para(tf, "•  " + text, size=17, colour=BLACK,
         first=text.startswith("That the index"), space_after=12)
footer(s, "Standing", num())

# ----------------------------------------------------------------- 26 output
s = slide()
action_title(s, "The work yields a dissertation and three papers", "books")
items = [("72 pp", "Dissertation", "Complete draft."),
         ("12 pp", "Paper 1", "The bound and the screening criterion. Physica A. "
                              "Independent of the scan."),
         ("12 pp", "Paper 2", "The instrument, the corpus, the field bound. "
                              "Physica A."),
         ("6 pp", "Paper 3", "The measured ceiling and the replication. Imaging "
                             "Neuroscience. One computation on its final run.")]
cw = 2.85
for i, (pages, name, body) in enumerate(items):
    x = M + i * (cw + 0.32)
    tf = textbox(s, x, 1.70, cw, 3.0)
    para(tf, pages, size=22, bold=True, colour=MIDBLUE, first=True, space_after=4)
    para(tf, name, size=20, bold=True, colour=BLACK, space_after=8)
    para(tf, body, size=15, colour=SLATE, space_after=0)
rule(s, M, 4.95, W - 2 * M)
tf = textbox(s, M, 5.20, W - 2 * M, 0.9)
para(tf, "Every number in all four comes from a re-runnable script. Nothing is "
         "transcribed by hand.", size=19, colour=BLACK, first=True, space_after=0)
footer(s, "Standing", num())

# ----------------------------------------------------------------- 27 conclusions
s = slide()
action_title(s, "Conclusions", "flag-banner")
tf = textbox(s, M, 1.50, W - 2 * M, 4.2)
conclusions = [
    ("A field observable can be measured from content. ",
     "The instrument runs, the corpus is complete, and reliability is measured "
     "rather than assumed: ICC = 0.8725."),
    ("It separates four corpora at 0.1068, and the separation is confounded "
     "with source. ", "Both halves are reported, the confound first."),
    ("The coupling is unidentified, so the thesis reports a bound instead. ",
     "With ΔX = 0.1241, α ≥ 4.294 at βJ = 2 — a constraint on any future "
     "proposal of this mechanism."),
    ("The instrument's own validity is now a measurable question. ",
     "The noise ceiling is 0.1517; one prediction pass remains."),
]
for i, (head, body) in enumerate(conclusions):
    rich(tf, [("%d.  " % (i + 1), True, MIDBLUE), (head, True, BLACK),
              (body, False, BLACK)], size=18, first=(i == 0), space_after=14)
rule(s, M, 5.85, W - 2 * M)
tf = textbox(s, M, 6.05, W - 2 * M, 0.5)
para(tf, "Brian Mwai  ·  1050555  ·  monarch-4iy.pages.dev  ·  "
         "github.com/brn-mwai/monarch", size=14, colour=SLATE, first=True,
     space_after=0)
footer(s, "Standing", num())

# ----------------------------------------------------------------- 28 references
s = slide()
action_title(s, "References")
refs = [
    "Ahmed, H., Traore, I., Saad, S. (2017). Detection of online fake news using "
    "n-gram analysis and machine learning. ISDDC.",
    "d'Ascoli, S. et al. (2025). TRIBE: a trimodal brain encoder for whole-brain "
    "fMRI response prediction. Meta FAIR. arXiv:2507.22229.",
    "Glasser, M. F. et al. (2016). A multi-modal parcellation of human cerebral "
    "cortex. Nature 536, 171–178.",
    "Kiesel, J. et al. (2019). SemEval-2019 Task 4: hyperpartisan news "
    "detection. SemEval.",
    "Korbel, J., Dahdoul, R., Thurner, S. (2026). Critical campaign spending in "
    "a double-random field Ising model of elections. Phys. Rev. Lett. 136, "
    "127402.",
    "Potthast, M. et al. (2018). The Webis clickbait challenge 2017. ECIR.",
    "Schaefer, A. et al. (2018). Local-global parcellation of the human cerebral "
    "cortex. Cereb. Cortex 28, 3095–3114.",
    "The Algonauts Project 2025 Challenge. arXiv:2501.00504.",
]
tf = textbox(s, M, 1.55, W - 2 * M, 5.0)
for i, ref in enumerate(refs):
    para(tf, ref, size=14, colour=BLACK, first=(i == 0), space_after=10)
footer(s, "Standing", num())

prs.save(OUT)
print("wrote", OUT)
print("slides", len(prs.slides.__iter__.__self__._sldIdLst))
